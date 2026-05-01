import fitz
from pptx import Presentation
import docx

def pdf_to_text(path: str) -> str:
    text = ""
    with fitz.open(path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def pptx_to_text(path: str) -> str:
    prs = Presentation(path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)

def docx_to_text(path: str) -> str:
    document = docx.Document(path)
    text = [para.text for para in document.paragraphs]
    return "\n".join(text)